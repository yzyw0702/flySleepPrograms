from pptx import Presentation


def testPptx():
    prs = Presentation()
    ttl_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(ttl_slide_layout)
    ttl = slide.shapes.title
    subtitle = slide.placeholders[1]

    ttl.text = 'Hello, World!'
    subtitle.text = 'python-pptx was here!'

    prs.save('test.pptx')

testPptx()